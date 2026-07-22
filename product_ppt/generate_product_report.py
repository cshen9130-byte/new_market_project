#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate private fund product introduction PPT and watermarked PDF (red style).

Usage:
    python generate_product_report.py

Outputs:
    output/私募基金产品介绍_红版_20260618.pptx
    output/私募基金产品介绍_红版_20260618.pdf
"""

from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

WORKSPACE = Path(__file__).resolve().parent
OUTPUT_DIR = WORKSPACE / "output"
END_DATE = pd.Timestamp("2026-06-18")
RISK_FREE_RATE = 0.02
WATERMARK_TEXT = "内部资料，请勿外传"
DISCLAIMER_TEXT = "注：历史业绩不代表未来表现，市场有风险，投资需谨慎。"
TRADING_DAYS_PER_YEAR = 252

PRODUCT_KEYWORDS = {
    "博孚利稳强8号FOF": "稳强8",
    "博孚利锡泰稳强FOF1号": "锡泰",
    "博孚利青璜2号": "青璜",
    "博孚利尊享15号FOF": "尊享15",
}

VOLATILITY_SECTIONS = ["低波", "中低波", "中波"]

SECTION_PRODUCT_ORDER: Dict[str, List[str]] = {
    "低波": ["博孚利稳强8号FOF", "博孚利锡泰稳强FOF1号"],
    "中低波": ["博孚利青璜2号"],
    "中波": ["博孚利尊享15号FOF"],
}

_SECTION_PRODUCT_ORDER: Dict[str, List[str]] = SECTION_PRODUCT_ORDER


def set_end_date(value: pd.Timestamp | str) -> None:
    global END_DATE
    END_DATE = pd.Timestamp(value)

@dataclass
class ReportTheme:
    color_primary: RGBColor
    color_accent: RGBColor
    color_text: RGBColor
    color_muted: RGBColor
    color_bg_light: RGBColor
    color_white: RGBColor
    color_chart_line: str
    color_chart_fill: str
    color_chart_fill_down: str
    color_up: RGBColor
    color_down: RGBColor
    section_styles: Dict[str, Dict[str, object]]
    cn_return_colors: bool = False


BLUE_THEME = ReportTheme(
    color_primary=RGBColor(0x1A, 0x36, 0x5D),
    color_accent=RGBColor(0xC9, 0xA2, 0x27),
    color_text=RGBColor(0x2D, 0x37, 0x48),
    color_muted=RGBColor(0x71, 0x85, 0x90),
    color_bg_light=RGBColor(0xF7, 0xFA, 0xFC),
    color_white=RGBColor(0xFF, 0xFF, 0xFF),
    color_chart_line="#1A365D",
    color_chart_fill="#E2E8F0",
    color_chart_fill_down="#E2E8F0",
    color_up=RGBColor(0x1A, 0x36, 0x5D),
    color_down=RGBColor(0x1A, 0x36, 0x5D),
    section_styles={
        "低波": {
            "accent": RGBColor(0x2B, 0x6C, 0xB6),
            "tagline": "波动可控 · 稳健收益",
        },
        "中低波": {
            "accent": RGBColor(0x38, 0x8E, 0x6C),
            "tagline": "适度波动 · 均衡配置",
        },
        "中波": {
            "accent": RGBColor(0xC9, 0x7A, 0x1E),
            "tagline": "弹性增强 · 追求超额",
        },
    },
)

RED_THEME = ReportTheme(
    color_primary=RGBColor(0x9B, 0x1C, 0x1C),
    color_accent=RGBColor(0xD4, 0xA0, 0x17),
    color_text=RGBColor(0x2D, 0x37, 0x48),
    color_muted=RGBColor(0x71, 0x85, 0x90),
    color_bg_light=RGBColor(0xFF, 0xF5, 0xF5),
    color_white=RGBColor(0xFF, 0xFF, 0xFF),
    color_chart_line="#C53030",
    color_chart_fill="#FED7D7",
    color_chart_fill_down="#C6F6D5",
    color_up=RGBColor(0xE5, 0x3E, 0x3E),
    color_down=RGBColor(0x38, 0xA1, 0x69),
    section_styles={
        "低波": {
            "accent": RGBColor(0xE5, 0x3E, 0x3E),
            "tagline": "波动可控 · 稳健收益",
        },
        "中低波": {
            "accent": RGBColor(0xC5, 0x30, 0x30),
            "tagline": "适度波动 · 均衡配置",
        },
        "中波": {
            "accent": RGBColor(0x9B, 0x2C, 0x2C),
            "tagline": "弹性增强 · 追求超额",
        },
    },
    cn_return_colors=True,
)

# Active theme (defaults to blue)
COLOR_PRIMARY = BLUE_THEME.color_primary
COLOR_ACCENT = BLUE_THEME.color_accent
COLOR_TEXT = BLUE_THEME.color_text
COLOR_MUTED = BLUE_THEME.color_muted
COLOR_BG_LIGHT = BLUE_THEME.color_bg_light
COLOR_WHITE = BLUE_THEME.color_white
COLOR_CHART_LINE = BLUE_THEME.color_chart_line
COLOR_CHART_FILL = BLUE_THEME.color_chart_fill
COLOR_CHART_FILL_DOWN = BLUE_THEME.color_chart_fill_down
COLOR_UP = BLUE_THEME.color_up
COLOR_DOWN = BLUE_THEME.color_down
SECTION_STYLES = BLUE_THEME.section_styles
USE_CN_RETURN_COLORS = BLUE_THEME.cn_return_colors


def apply_theme(theme: ReportTheme) -> None:
    global COLOR_PRIMARY, COLOR_ACCENT, COLOR_TEXT, COLOR_MUTED
    global COLOR_BG_LIGHT, COLOR_WHITE, COLOR_CHART_LINE, COLOR_CHART_FILL
    global COLOR_CHART_FILL_DOWN, COLOR_UP, COLOR_DOWN, SECTION_STYLES
    global USE_CN_RETURN_COLORS

    COLOR_PRIMARY = theme.color_primary
    COLOR_ACCENT = theme.color_accent
    COLOR_TEXT = theme.color_text
    COLOR_MUTED = theme.color_muted
    COLOR_BG_LIGHT = theme.color_bg_light
    COLOR_WHITE = theme.color_white
    COLOR_CHART_LINE = theme.color_chart_line
    COLOR_CHART_FILL = theme.color_chart_fill
    COLOR_CHART_FILL_DOWN = theme.color_chart_fill_down
    COLOR_UP = theme.color_up
    COLOR_DOWN = theme.color_down
    SECTION_STYLES = theme.section_styles
    USE_CN_RETURN_COLORS = theme.cn_return_colors

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------

@dataclass
class ProductInfo:
    name: str
    volatility_type: str = ""
    filing_no: str = ""
    manager: str = ""
    inception_date: str = ""
    aum: str = ""
    fund_manager: str = ""


@dataclass
class PerformanceMetrics:
    unit_nav: float = 0.0
    total_return: float = 0.0
    ytd_return: float = 0.0
    annualized_return: float = 0.0
    max_drawdown: float = 0.0
    sharpe_ratio: float = 0.0


@dataclass
class IntervalStats:
    period_return: float = 0.0
    annualized_return: float = 0.0
    annualized_volatility: float = 0.0
    sharpe_ratio: float = 0.0
    calmar_ratio: float = 0.0
    downside_risk: float = 0.0
    max_drawdown: float = 0.0
    max_dd_recovery_days: int = 0
    max_non_high_days: int = 0


@dataclass
class ProductReport:
    info: ProductInfo
    nav_df: pd.DataFrame
    basic: PerformanceMetrics = field(default_factory=PerformanceMetrics)
    interval_inception: IntervalStats = field(default_factory=IntervalStats)
    interval_3y: IntervalStats = field(default_factory=IntervalStats)
    interval_3y_start: pd.Timestamp = field(default_factory=lambda: END_DATE)
    monthly_table: pd.DataFrame = field(default_factory=pd.DataFrame)


# ---------------------------------------------------------------------------
# Data loading
# ---------------------------------------------------------------------------

def _extract_field_value(line: str, label: str) -> Optional[str]:
    match = re.match(rf"^{re.escape(label)}\s*[：:]\s*(.+)$", line)
    return match.group(1).strip() if match else None


def parse_basic_information(path: Path) -> Dict[str, ProductInfo]:
    text = path.read_text(encoding="utf-8")
    blocks = re.split(r"\n{3,}", text.strip())
    products: Dict[str, ProductInfo] = {}
    field_map = {
        "波动类型": "volatility_type",
        "备案编号": "filing_no",
        "私募管理人": "manager",
        "产品成立时间": "inception_date",
        "公司管理规模": "aum",
        "基金经理": "fund_manager",
    }

    for block in blocks:
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
        name = lines[0]
        info = ProductInfo(name=name)
        for line in lines[1:]:
            for label, attr in field_map.items():
                value = _extract_field_value(line, label)
                if value is not None:
                    setattr(info, attr, value)
                    break
        products[name] = info
    return products


def find_nav_file(name: str, directory: Path) -> Path:
    keyword = PRODUCT_KEYWORDS[name]
    matches = [
        p for p in directory.iterdir()
        if p.is_file() and p.suffix.lower() in {".xlsx", ".csv"} and keyword in p.name
    ]
    if not matches:
        raise FileNotFoundError(f"No NAV file found for {name} (keyword: {keyword})")
    # Pick the file whose name sorts last (date suffix like 20260625 > 20260610)
    return sorted(matches, key=lambda p: p.name)[-1]


def resolve_nav_path(name: str, directory: Path, nav_file: Optional[str] = None) -> Path:
    if nav_file:
        path = directory / nav_file
        if not path.is_file():
            raise FileNotFoundError(f"NAV file not found for {name}: {path}")
        return path
    return find_nav_file(name, directory)


def load_nav_data(path: Path) -> pd.DataFrame:
    if path.suffix.lower() in {".csv", ".txt"}:
        df = pd.read_csv(path, encoding="utf-8-sig")
    else:
        df = pd.read_excel(path, sheet_name=0, header=0)
    expected_cols = ["date", "unit_nav", "cum_nav", "adj_nav", "pct_change", "benchmark"]
    if len(df.columns) > len(expected_cols):
        df = df.iloc[:, : len(expected_cols)]
    df.columns = expected_cols
    df["date"] = pd.to_datetime(df["date"])
    for col in ["unit_nav", "cum_nav", "adj_nav", "benchmark"]:
        df[col] = pd.to_numeric(df[col], errors="coerce")
    pct = df["pct_change"].astype(str).str.replace("%", "", regex=False)
    df["period_return"] = pd.to_numeric(pct, errors="coerce") / 100

    nav_col = "adj_nav" if df["adj_nav"].notna().sum() > 0 else "cum_nav"
    if df[nav_col].isna().all():
        nav_col = "unit_nav"
    df["nav"] = df[nav_col]
    df = df.dropna(subset=["date", "nav"]).sort_values("date").reset_index(drop=True)
    df = df[df["date"] <= END_DATE].copy()
    return df


# ---------------------------------------------------------------------------
# Performance calculations
# ---------------------------------------------------------------------------

def _annualization_factor(dates: pd.Series) -> float:
    if len(dates) < 2:
        return 52.0
    gaps = dates.diff().dt.days.dropna()
    median_gap = gaps.median() if len(gaps) else 7
    if median_gap <= 0:
        median_gap = 7
    periods_per_year = 365.0 / median_gap
    return periods_per_year


def _period_returns(df: pd.DataFrame) -> pd.Series:
    if df["period_return"].notna().sum() >= len(df) * 0.5:
        returns = df["period_return"].dropna()
    else:
        returns = df["nav"].pct_change().dropna()
    return returns


def _total_return(nav_start: float, nav_end: float) -> float:
    if nav_start <= 0:
        return 0.0
    return nav_end / nav_start - 1


def _annualized_return(total_ret: float, days: int) -> float:
    if days <= 0:
        return 0.0
    years = days / 365.25
    if years <= 0:
        return 0.0
    return (1 + total_ret) ** (1 / years) - 1


def _max_drawdown(nav: pd.Series) -> Tuple[float, int, int]:
    """Return max drawdown (negative), recovery days, longest non-high streak."""
    if nav.empty:
        return 0.0, 0, 0

    peak = nav.cummax()
    drawdown = nav / peak - 1
    max_dd = drawdown.min()

    trough_date = drawdown.idxmin()
    trough_pos = nav.index.get_loc(trough_date)
    peak_before = peak.loc[trough_date]
    recovery_days = -1
    for later_date in nav.index[trough_pos + 1 :]:
        if nav.loc[later_date] >= peak_before:
            recovery_days = (later_date - trough_date).days
            break

    dates = nav.index
    running_peak = nav.iloc[0]
    streak_start = dates[0]
    max_streak = 0
    for i in range(1, len(nav)):
        if nav.iloc[i] >= running_peak:
            running_peak = nav.iloc[i]
            max_streak = max(max_streak, (dates[i] - streak_start).days)
            streak_start = dates[i]
    if len(nav) > 1:
        max_streak = max(max_streak, (dates[-1] - streak_start).days)

    return max_dd, recovery_days, max_streak


def _downside_risk(returns: pd.Series, periods_per_year: float) -> float:
    negative = returns[returns < 0]
    if negative.empty:
        return 0.0
    return negative.std() * np.sqrt(periods_per_year)


def _sharpe(annual_return: float, annual_vol: float, rf: float = RISK_FREE_RATE) -> float:
    if annual_vol <= 0:
        return 0.0
    return (annual_return - rf) / annual_vol


def calc_basic_performance(df: pd.DataFrame) -> PerformanceMetrics:
    df = df.copy()
    df = df.set_index("date")
    nav = df["nav"]
    returns = _period_returns(df.reset_index())

    nav_end = nav.iloc[-1]
    nav_start = nav.iloc[0]
    total_ret = _total_return(nav_start, nav_end)
    days = (nav.index[-1] - nav.index[0]).days
    ann_ret = _annualized_return(total_ret, days)

    year_start = pd.Timestamp(f"{END_DATE.year}-01-01")
    ytd_df = df[df.index >= year_start]
    prev_year = df[df.index < year_start]
    # Use last nav of previous year as base so YTD matches the monthly table's compounded 全年
    if not prev_year.empty:
        base_nav = prev_year["nav"].iloc[-1]
        ytd_ret = _total_return(base_nav, nav_end)
    elif not ytd_df.empty:
        ytd_ret = _total_return(ytd_df["nav"].iloc[0], nav_end)
    else:
        ytd_ret = 0.0

    max_dd, _, _ = _max_drawdown(nav)
    periods_per_year = _annualization_factor(df.reset_index()["date"])
    vol = returns.std() * np.sqrt(periods_per_year) if len(returns) > 1 else 0.0
    sharpe = _sharpe(ann_ret, vol)

    raw_unit_nav = df["unit_nav"].iloc[-1] if "unit_nav" in df.columns and df["unit_nav"].notna().any() else nav_end

    return PerformanceMetrics(
        unit_nav=raw_unit_nav,
        total_return=total_ret,
        ytd_return=ytd_ret,
        annualized_return=ann_ret,
        max_drawdown=max_dd,
        sharpe_ratio=sharpe,
    )


def _filter_last_n_years(df: pd.DataFrame, years: int = 3) -> pd.DataFrame:
    start = END_DATE - pd.DateOffset(years=years)
    subset = df[df["date"] >= start].copy()
    return subset if len(subset) >= 2 else df.copy()


def calc_interval_stats(df: pd.DataFrame) -> IntervalStats:
    df = df.copy().set_index("date")
    nav = df["nav"]
    returns = _period_returns(df.reset_index())
    periods_per_year = _annualization_factor(df.reset_index()["date"])

    total_ret = _total_return(nav.iloc[0], nav.iloc[-1])
    days = (nav.index[-1] - nav.index[0]).days
    ann_ret = _annualized_return(total_ret, days)
    ann_vol = returns.std() * np.sqrt(periods_per_year) if len(returns) > 1 else 0.0
    max_dd, recovery_days, max_non_high = _max_drawdown(nav)
    downside = _downside_risk(returns, periods_per_year)
    sharpe = _sharpe(ann_ret, ann_vol)
    calmar = ann_ret / abs(max_dd) if max_dd < 0 else 0.0

    return IntervalStats(
        period_return=total_ret,
        annualized_return=ann_ret,
        annualized_volatility=ann_vol,
        sharpe_ratio=sharpe,
        calmar_ratio=calmar,
        downside_risk=downside,
        max_drawdown=max_dd,
        max_dd_recovery_days=recovery_days,
        max_non_high_days=max_non_high,
    )


def calc_monthly_returns(df: pd.DataFrame, product_name: str) -> pd.DataFrame:
    df = df.copy().set_index("date")
    month_end = df["nav"].resample("ME").last().dropna()
    monthly_ret = month_end.pct_change()
    if not month_end.empty:
        monthly_ret.iloc[0] = month_end.iloc[0] / df["nav"].iloc[0] - 1

    records = []
    for dt, ret in monthly_ret.items():
        records.append({"year": dt.year, "month": dt.month, "return": ret})

    if not records:
        return pd.DataFrame()

    ret_df = pd.DataFrame(records)
    years = sorted(ret_df["year"].unique())
    rows = []
    for year in years:
        row = {"年份": str(year), "基金名称": product_name}
        year_data = ret_df[ret_df["year"] == year]
        annual = 1.0
        for m in range(1, 13):
            sub = year_data[year_data["month"] == m]
            if sub.empty:
                row[f"{m}月"] = ""
            else:
                val = sub["return"].iloc[0]
                row[f"{m}月"] = _fmt_pct(val)
                annual *= 1 + val
        row["全年"] = _fmt_pct(annual - 1)
        rows.append(row)

    columns = ["年份", "基金名称"] + [f"{m}月" for m in range(1, 13)] + ["全年"]
    return pd.DataFrame(rows, columns=columns)


def _build_product_report(name: str, info: ProductInfo, nav_df: pd.DataFrame) -> ProductReport:
    nav_3y = _filter_last_n_years(nav_df, years=3)
    return ProductReport(
        info=info,
        nav_df=nav_df,
        basic=calc_basic_performance(nav_df),
        interval_inception=calc_interval_stats(nav_df),
        interval_3y=calc_interval_stats(nav_3y),
        interval_3y_start=nav_3y["date"].min(),
        monthly_table=calc_monthly_returns(nav_df, name),
    )


def build_reports(workspace: Path, config: Optional[dict] = None) -> List[ProductReport]:
    global _SECTION_PRODUCT_ORDER

    if config is None:
        config_path = workspace / "report_config.json"
        if config_path.is_file():
            config = json.loads(config_path.read_text(encoding="utf-8"))

    if config:
        set_end_date(config["end_date"])
        _SECTION_PRODUCT_ORDER = config.get("section_order", {})
        products_cfg: Dict[str, dict] = config.get("products", {})
        ordered_names: List[str] = []
        for section in VOLATILITY_SECTIONS:
            for name in _SECTION_PRODUCT_ORDER.get(section, []):
                if name not in ordered_names:
                    ordered_names.append(name)
        for name in products_cfg:
            if name not in ordered_names:
                ordered_names.append(name)

        reports: List[ProductReport] = []
        for name in ordered_names:
            product_cfg = products_cfg.get(name, {})
            info = ProductInfo(
                name=name,
                volatility_type=product_cfg.get("volatility_type", ""),
                filing_no=product_cfg.get("filing_no", ""),
                manager=product_cfg.get("manager", ""),
                inception_date=product_cfg.get("inception_date", ""),
                aum=product_cfg.get("aum", ""),
                fund_manager=product_cfg.get("fund_manager", ""),
            )
            nav_path = resolve_nav_path(name, workspace, product_cfg.get("nav_file"))
            nav_df = load_nav_data(nav_path)
            reports.append(_build_product_report(name, info, nav_df))
        return reports

    _SECTION_PRODUCT_ORDER = SECTION_PRODUCT_ORDER
    basic_path = workspace / "basic_information.txt"
    products = parse_basic_information(basic_path)
    reports = []

    for name in PRODUCT_KEYWORDS:
        info = products.get(name, ProductInfo(name=name))
        nav_path = find_nav_file(name, workspace)
        nav_df = load_nav_data(nav_path)
        reports.append(_build_product_report(name, info, nav_df))
    return reports


def group_reports_by_section(reports: List[ProductReport]) -> List[Tuple[str, List[ProductReport]]]:
    by_name = {r.info.name: r for r in reports}
    grouped: List[Tuple[str, List[ProductReport]]] = []
    for section in VOLATILITY_SECTIONS:
        section_reports: List[ProductReport] = []
        for name in _SECTION_PRODUCT_ORDER.get(section, []):
            if name in by_name:
                section_reports.append(by_name[name])
        for r in reports:
            if r.info.volatility_type == section and r not in section_reports:
                section_reports.append(r)
        if section_reports:
            grouped.append((section, section_reports))
    return grouped


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------

def _fmt_pct(value: Optional[float], digits: int = 2) -> str:
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return "--"
    return f"{value * 100:.{digits}f}%"


def _fmt_num(value: Optional[float], digits: int = 4) -> str:
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return "--"
    return f"{value:.{digits}f}"


def _fmt_days(value: int) -> str:
    if value < 0:
        return "未回补"
    return f"{value}"


def _pct_text_color(text: str, default: Optional[RGBColor] = None) -> RGBColor:
    fallback = default or COLOR_TEXT
    if not USE_CN_RETURN_COLORS or not text or text == "--" or not text.endswith("%"):
        return fallback
    try:
        value = float(text.replace("%", ""))
    except ValueError:
        return fallback
    if value > 0:
        return COLOR_UP
    if value < 0:
        return COLOR_DOWN
    return fallback


# ---------------------------------------------------------------------------
# Chart generation
# ---------------------------------------------------------------------------

def create_nav_chart(df: pd.DataFrame, title: str, output_path: Path) -> None:
    df = df.copy()
    base_nav = df["nav"].iloc[0]
    df["return_pct"] = (df["nav"] / base_nav - 1) * 100

    fig, ax = plt.subplots(figsize=(12.5, 3.35), dpi=120)
    if USE_CN_RETURN_COLORS:
        ax.fill_between(
            df["date"],
            df["return_pct"],
            0,
            where=(df["return_pct"] >= 0),
            color=COLOR_CHART_FILL,
            alpha=0.75,
            interpolate=True,
        )
        ax.fill_between(
            df["date"],
            df["return_pct"],
            0,
            where=(df["return_pct"] < 0),
            color=COLOR_CHART_FILL_DOWN,
            alpha=0.75,
            interpolate=True,
        )
    else:
        ax.fill_between(df["date"], df["return_pct"], 0, color=COLOR_CHART_FILL, alpha=0.6)
    ax.plot(df["date"], df["return_pct"], color=COLOR_CHART_LINE, linewidth=2.2, label="累计收益率")
    ax.axhline(0, color="#A0AEC0", linewidth=0.8, linestyle="--")

    ax.set_title(title, fontsize=12, fontweight="bold", color=COLOR_CHART_LINE, pad=8)
    ax.set_ylabel("收益率 (%)", fontsize=9)
    ax.set_xlabel("")
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m"))
    ax.xaxis.set_major_locator(mdates.YearLocator())
    ax.grid(True, linestyle=":", alpha=0.45)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=7)
    fig.subplots_adjust(left=0.06, right=0.99, top=0.88, bottom=0.18)
    fig.savefig(output_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ---------------------------------------------------------------------------
# PPT helpers
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(
    slide,
    title: str,
    subtitle: str = "",
    accent_color: Optional[RGBColor] = None,
) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0.82), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color or COLOR_ACCENT
    accent.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(10), Inches(0.5))
    p = txb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.48), Inches(10), Inches(0.3))
        p2 = txb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)


SECTION_LABEL_HEIGHT = 0.32


def _add_section_label(slide, text: str, left: float, top: float, width: float = 3.0) -> float:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(SECTION_LABEL_HEIGHT))
    tf = box.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    return top + SECTION_LABEL_HEIGHT


def _add_info_card(slide, info: ProductInfo, left: float, top: float, width: float, height: float) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG_LIGHT
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    items = [
        ("波动类型", info.volatility_type),
        ("备案编号", info.filing_no),
        ("私募管理人", info.manager),
        ("成立时间", info.inception_date),
        ("管理规模", info.aum),
        ("基金经理", info.fund_manager),
    ]
    txb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(width - 0.3), Inches(height - 0.2))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (label, value) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{label}：{value}"
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(3)


def _add_metric_cards(slide, metrics: List[Tuple[str, str]], left: float, top: float, cols: int = 3) -> None:
    card_w, card_h = 2.05, 0.72
    gap_x, gap_y = 0.12, 0.1
    for i, (label, value) in enumerate(metrics):
        row, col = divmod(i, cols)
        x = left + col * (card_w + gap_x)
        y = top + row * (card_h + gap_y)

        card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

        lbl = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08), Inches(card_w - 0.2), Inches(0.25))
        lp = lbl.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(8)
        lp.font.color.rgb = COLOR_MUTED

        val = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.3), Inches(card_w - 0.2), Inches(0.35))
        vp = val.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(14)
        vp.font.bold = True
        vp.font.color.rgb = _pct_text_color(value, COLOR_PRIMARY)


def _add_dual_stats_table(
    slide,
    inception: IntervalStats,
    last_3y: IntervalStats,
    left: float,
    top: float,
    width: float,
    height: float = 2.65,
) -> float:
    rows = [
        ("区间收益", _fmt_pct(inception.period_return), _fmt_pct(last_3y.period_return)),
        ("年化收益", _fmt_pct(inception.annualized_return), _fmt_pct(last_3y.annualized_return)),
        ("年化波动率", _fmt_pct(inception.annualized_volatility), _fmt_pct(last_3y.annualized_volatility)),
        ("夏普比率 (Rf=2.00%)", _fmt_num(inception.sharpe_ratio, 2), _fmt_num(last_3y.sharpe_ratio, 2)),
        ("卡玛比率", _fmt_num(inception.calmar_ratio, 2), _fmt_num(last_3y.calmar_ratio, 2)),
        ("下行风险", _fmt_pct(inception.downside_risk), _fmt_pct(last_3y.downside_risk)),
        ("最大回撤", _fmt_pct(inception.max_drawdown), _fmt_pct(last_3y.max_drawdown)),
        ("最大回撤回补期 (天)", _fmt_days(inception.max_dd_recovery_days), _fmt_days(last_3y.max_dd_recovery_days)),
        (
            "最长连续不创新高 (天)",
            str(inception.max_non_high_days),
            str(last_3y.max_non_high_days),
        ),
    ]

    n_rows = len(rows) + 1
    n_cols = 3
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    table.columns[0].width = Inches(width * 0.42)
    table.columns[1].width = Inches(width * 0.29)
    table.columns[2].width = Inches(width * 0.29)
    row_height = Inches(height / n_rows)
    for row in table.rows:
        row.height = row_height

    headers = ["指标", "成立以来", "最近三年"]
    for j, title in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = title
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(1)
        cell.margin_bottom = Pt(1)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER

    for r, (label, v_inception, v_3y) in enumerate(rows, start=1):
        values = [label, v_inception, v_3y]
        for j, value in enumerate(values):
            cell = table.rows[r].cells[j]
            cell.text = value
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BG_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9 if j > 0 else 8)
                p.font.bold = j > 0
                p.font.color.rgb = _pct_text_color(value) if j > 0 else COLOR_TEXT
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

    return top + height


def _add_monthly_table(
    slide,
    monthly: pd.DataFrame,
    left: float,
    top: float,
    width: float,
    max_height: Optional[float] = None,
) -> float:
    if monthly.empty:
        return top

    cols = [c for c in monthly.columns if c != "基金名称"]
    n_rows = len(monthly) + 1
    n_cols = len(cols)
    height_cap = max_height if max_height is not None else 2.45
    height = min(0.34 * n_rows, height_cap)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    row_height = Inches(height / n_rows)
    for row in table.rows:
        row.height = row_height

    for j, col_name in enumerate(cols):
        cell = table.rows[0].cells[j]
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER

    for i, row in monthly.iterrows():
        for j, col_name in enumerate(cols):
            cell = table.rows[i + 1].cells[j]
            val = str(row[col_name]) if pd.notna(row[col_name]) else ""
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BG_LIGHT
            is_return_col = col_name not in ("年份", "基金名称")
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9 if is_return_col else 8)
                p.font.bold = is_return_col
                p.font.color.rgb = _pct_text_color(val) if is_return_col else COLOR_TEXT
                p.alignment = PP_ALIGN.CENTER

    return top + height


def _add_footer(slide, top: float = 7.05, *, include_disclaimer: bool = False) -> None:
    if include_disclaimer:
        disclaimer_box = slide.shapes.add_textbox(Inches(0.4), Inches(top), Inches(9.6), Inches(0.3))
        dp = disclaimer_box.text_frame.paragraphs[0]
        dp.text = DISCLAIMER_TEXT
        dp.font.size = Pt(8)
        dp.font.color.rgb = COLOR_MUTED
        dp.alignment = PP_ALIGN.LEFT

    watermark_left = 9.8 if include_disclaimer else 0.5
    watermark_width = 3.1 if include_disclaimer else 12.0
    watermark_box = slide.shapes.add_textbox(
        Inches(watermark_left), Inches(top), Inches(watermark_width), Inches(0.3)
    )
    wp = watermark_box.text_frame.paragraphs[0]
    wp.text = WATERMARK_TEXT
    wp.font.size = Pt(8)
    wp.font.color.rgb = COLOR_MUTED
    wp.alignment = PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# PPT generation
# ---------------------------------------------------------------------------

def create_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.2), Inches(13.333), Inches(2.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.2))
    p = txb.text_frame.paragraphs[0]
    p.text = "私募产品历史业绩"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(0.5))
    sp = sub.text_frame.paragraphs[0]
    sp.text = f"数据截至 {END_DATE.strftime('%Y年%m月%d日')}  |  低波 · 中低波 · 中波  |  内部资料"
    sp.font.size = Pt(14)
    sp.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    sp.alignment = PP_ALIGN.CENTER

    _add_footer(slide, include_disclaimer=True)


def create_contents_slide(prs: Presentation, sections: List[Tuple[str, List[ProductReport]]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)
    _add_header_bar(slide, "产品目录", f"数据截至 {END_DATE.strftime('%Y年%m月%d日')}")

    card_top = 1.25
    card_h = 1.75
    card_gap = 0.22
    card_w = 3.9
    starts = [0.45, 4.72, 8.98]

    for idx, (section, section_reports) in enumerate(sections[:3]):
        style = SECTION_STYLES.get(section, {})
        accent = style.get("accent", COLOR_ACCENT)
        left = starts[idx] if idx < len(starts) else 0.45 + idx * (card_w + card_gap)

        card = slide.shapes.add_shape(1, Inches(left), Inches(card_top), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_LIGHT
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

        stripe = slide.shapes.add_shape(1, Inches(left), Inches(card_top), Inches(card_w), Inches(0.12))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent  # type: ignore[arg-type]
        stripe.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(left + 0.18), Inches(card_top + 0.22), Inches(card_w - 0.36), Inches(0.45)
        )
        tp = title_box.text_frame.paragraphs[0]
        tp.text = section
        tp.font.size = Pt(20)
        tp.font.bold = True
        tp.font.color.rgb = COLOR_PRIMARY

        tagline = str(style.get("tagline", ""))
        if tagline:
            tag_box = slide.shapes.add_textbox(
                Inches(left + 0.18), Inches(card_top + 0.62), Inches(card_w - 0.36), Inches(0.28)
            )
            tgp = tag_box.text_frame.paragraphs[0]
            tgp.text = tagline
            tgp.font.size = Pt(8)
            tgp.font.color.rgb = COLOR_MUTED

        list_box = slide.shapes.add_textbox(
            Inches(left + 0.18), Inches(card_top + 0.95), Inches(card_w - 0.36), Inches(0.7)
        )
        tf = list_box.text_frame
        tf.word_wrap = True
        for i, report in enumerate(section_reports):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"· {report.info.name}"
            p.font.size = Pt(8)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(2)

    _add_footer(slide)


def create_section_divider_slide(
    prs: Presentation,
    section: str,
    section_reports: List[ProductReport],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)

    style = SECTION_STYLES.get(section, {})
    accent = style.get("accent", COLOR_ACCENT)
    tagline = str(style.get("tagline", ""))

    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.45), Inches(13.333), Inches(2.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(4.98), Inches(13.333), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent  # type: ignore[arg-type]
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.95), Inches(11), Inches(1.0))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = f"{section}产品"
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_WHITE
    tp.alignment = PP_ALIGN.CENTER

    if tagline:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.85), Inches(11), Inches(0.45))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = tagline
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        sp.alignment = PP_ALIGN.CENTER

    list_box = slide.shapes.add_textbox(Inches(3.2), Inches(5.35), Inches(7), Inches(1.2))
    tf = list_box.text_frame
    intro = tf.paragraphs[0]
    intro.text = "本章节产品"
    intro.font.size = Pt(11)
    intro.font.bold = True
    intro.font.color.rgb = COLOR_PRIMARY
    intro.alignment = PP_ALIGN.CENTER

    for i, report in enumerate(section_reports):
        p = tf.add_paragraph()
        p.text = report.info.name
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)

    _add_footer(slide)


def create_product_overview_slide(prs: Presentation, report: ProductReport, chart_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)

    period = f"{report.nav_df['date'].min().strftime('%Y-%m-%d')} 至 {END_DATE.strftime('%Y-%m-%d')}"
    vtype = report.info.volatility_type
    section_style = SECTION_STYLES.get(vtype, {})
    accent = section_style.get("accent")
    subtitle = f"{vtype}  |  区间：{period}" if vtype else f"区间：{period}"
    _add_header_bar(slide, report.info.name, subtitle, accent_color=accent)  # type: ignore[arg-type]

    _add_section_label(slide, "基本信息", 0.4, 1.0)
    _add_info_card(slide, report.info, 0.4, 1.3, 3.5, 1.55)

    _add_section_label(slide, "基础业绩", 4.2, 1.0)
    basic_metrics = [
        ("单位净值", _fmt_num(report.basic.unit_nav)),
        ("成立以来收益", _fmt_pct(report.basic.total_return)),
        ("今年以来收益", _fmt_pct(report.basic.ytd_return)),
        ("成立以来年化", _fmt_pct(report.basic.annualized_return)),
        ("成立以来最大回撤", _fmt_pct(report.basic.max_drawdown)),
        ("成立以来夏普比率", _fmt_num(report.basic.sharpe_ratio, 2)),
    ]
    _add_metric_cards(slide, basic_metrics, 4.2, 1.3, cols=3)

    _add_section_label(slide, "净值走势（累计收益率）", 0.4, 2.95)
    slide.shapes.add_picture(
        str(chart_path),
        Inches(0.4),
        Inches(3.2),
        width=Inches(12.5),
        height=Inches(3.35),
    )

    _add_footer(slide, include_disclaimer=True)


def create_product_stats_slide(prs: Presentation, report: ProductReport) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_WHITE)

    inception_start = report.nav_df["date"].min().strftime("%Y-%m-%d")
    last_3y_start = report.interval_3y_start.strftime("%Y-%m-%d")
    end_str = END_DATE.strftime("%Y-%m-%d")
    subtitle = (
        f"成立以来：{inception_start} 至 {end_str}  |  "
        f"最近三年：{last_3y_start} 至 {end_str}"
    )
    vtype = report.info.volatility_type
    section_style = SECTION_STYLES.get(vtype, {})
    accent = section_style.get("accent")
    header_subtitle = f"{vtype}  |  {subtitle}" if vtype else subtitle
    _add_header_bar(slide, report.info.name, header_subtitle, accent_color=accent)  # type: ignore[arg-type]

    label_table_gap = 0.12
    section_gap = 0.38
    note_top = 6.7
    note_gap = 0.12

    stats_label_top = 1.0
    stats_table_top = _add_section_label(slide, "区间统计", 0.4, stats_label_top, width=4.0) + label_table_gap
    stats_table_height = 2.05
    stats_table_bottom = _add_dual_stats_table(
        slide,
        report.interval_inception,
        report.interval_3y,
        0.4,
        stats_table_top,
        12.5,
        height=stats_table_height,
    )
    monthly_label_top = stats_table_bottom + section_gap
    monthly_table_top = _add_section_label(slide, "区间收益（分月）", 0.4, monthly_label_top, width=4.0) + label_table_gap
    monthly_max_height = note_top - monthly_table_top - note_gap
    _add_monthly_table(
        slide,
        report.monthly_table,
        0.4,
        monthly_table_top,
        12.5,
        max_height=monthly_max_height,
    )

    note = slide.shapes.add_textbox(Inches(0.4), Inches(note_top), Inches(9.6), Inches(0.45))
    tf = note.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    note_lines = [
        "注：区间统计基于复权净值，夏普比率无风险利率 Rf=2.00%，最近三年为截至日向前滚动三年，月度收益为自然月收益率。",
        DISCLAIMER_TEXT,
    ]
    for i, line in enumerate(note_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(7)
        p.font.color.rgb = COLOR_MUTED
        p.space_before = Pt(0)
        p.space_after = Pt(1)

    # Watermark on the right of the same row
    wm = slide.shapes.add_textbox(Inches(9.8), Inches(note_top + 0.1), Inches(3.1), Inches(0.28))
    wp = wm.text_frame.paragraphs[0]
    wp.text = WATERMARK_TEXT
    wp.font.size = Pt(8)
    wp.font.color.rgb = COLOR_MUTED
    wp.alignment = PP_ALIGN.RIGHT


def generate_ppt(reports: List[ProductReport], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    sections = group_reports_by_section(reports)

    create_title_slide(prs)
    create_contents_slide(prs, sections)

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        for section, section_reports in sections:
            create_section_divider_slide(prs, section, section_reports)
            for report in section_reports:
                chart_path = tmp / f"{report.info.name}_chart.png"
                create_nav_chart(
                    report.nav_df,
                    f"{report.info.name} 累计收益率走势",
                    chart_path,
                )
                create_product_overview_slide(prs, report, chart_path)
                create_product_stats_slide(prs, report)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))


# ---------------------------------------------------------------------------
# PDF conversion & watermark
# ---------------------------------------------------------------------------

def _find_soffice() -> str | None:
    candidates: list[str] = []
    env_path = os.environ.get("SOFFICE_PATH", "").strip()
    if env_path:
        candidates.append(env_path)
    for cmd in ("soffice", "libreoffice"):
        found = shutil.which(cmd)
        if found:
            candidates.append(found)
    candidates.extend(
        [
            "/usr/bin/soffice",
            "/usr/local/bin/soffice",
            "/usr/bin/libreoffice",
            "/usr/lib/libreoffice/program/soffice",
        ]
    )

    seen: set[str] = set()
    for candidate in candidates:
        if not candidate or candidate in seen:
            continue
        seen.add(candidate)
        if os.path.isfile(candidate) and os.access(candidate, os.X_OK):
            return candidate
    return None


def _pptx_to_pdf_powerpoint(pptx_path: Path, pdf_path: Path) -> None:
    """Convert PPTX to PDF via PowerPoint COM (Windows)."""
    import comtypes.client

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    try:
        presentation = powerpoint.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        presentation.SaveAs(str(pdf_path.resolve()), 32)  # ppSaveAsPDF
        presentation.Close()
    finally:
        powerpoint.Quit()


def _pptx_to_pdf_soffice(soffice: str, pptx_path: Path, pdf_path: Path) -> None:
    """Convert PPTX to PDF via LibreOffice headless (Linux/macOS).

    Uses ASCII temp filenames — Chinese paths often break soffice on Linux servers.
    """
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    work_dir = Path(tempfile.mkdtemp(prefix="lo_pptx2pdf_"))
    lo_profile = work_dir / "profile"
    lo_profile.mkdir(parents=True, exist_ok=True)
    ascii_pptx = work_dir / "report.pptx"
    ascii_pdf = work_dir / "report.pdf"

    env = os.environ.copy()
    env["HOME"] = str(work_dir)
    env.setdefault("LANG", "C.UTF-8")
    env.setdefault("LC_ALL", "C.UTF-8")
    env["SAL_USE_VCLPLUGIN"] = "svp"

    try:
        shutil.copy2(pptx_path, ascii_pptx)
        profile_uri = lo_profile.resolve().as_uri()
        cmd = [
            soffice,
            f"-env:UserInstallation={profile_uri}",
            "--headless",
            "--nologo",
            "--nolockcheck",
            "--nodefault",
            "--norestore",
            "--invisible",
            "--convert-to",
            "pdf:impress_pdf_Export",
            "--outdir",
            str(work_dir),
            str(ascii_pptx),
        ]
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=300,
            check=False,
            env=env,
        )
        detail = (result.stderr or result.stdout or "").strip()

        produced = ascii_pdf if ascii_pdf.is_file() else None
        if produced is None:
            pdfs = sorted(work_dir.glob("*.pdf"), key=lambda p: p.stat().st_mtime, reverse=True)
            produced = pdfs[0] if pdfs else None

        if produced is None:
            # Retry once with the simpler filter name (some installs lack impress_pdf_Export).
            if "impress_pdf_Export" in " ".join(cmd):
                cmd_retry = [
                    soffice,
                    f"-env:UserInstallation={profile_uri}",
                    "--headless",
                    "--nologo",
                    "--nolockcheck",
                    "--nodefault",
                    "--norestore",
                    "--invisible",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(work_dir),
                    str(ascii_pptx),
                ]
                result = subprocess.run(
                    cmd_retry,
                    capture_output=True,
                    text=True,
                    timeout=300,
                    check=False,
                    env=env,
                )
                detail = (result.stderr or result.stdout or "").strip()
                if ascii_pdf.is_file():
                    produced = ascii_pdf
                else:
                    pdfs = sorted(work_dir.glob("*.pdf"), key=lambda p: p.stat().st_mtime, reverse=True)
                    produced = pdfs[0] if pdfs else None

        if produced is None:
            raise RuntimeError(
                f"LibreOffice PDF conversion failed (exit={result.returncode}): {detail or 'no PDF produced'}"
            )

        if pdf_path.exists():
            pdf_path.unlink()
        shutil.copy2(produced, pdf_path)
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


def pptx_to_pdf(pptx_path: Path, pdf_path: Path) -> None:
    if sys.platform == "win32":
        _pptx_to_pdf_powerpoint(pptx_path, pdf_path)
        return

    soffice = _find_soffice()
    if soffice:
        _pptx_to_pdf_soffice(soffice, pptx_path, pdf_path)
        return

    raise RuntimeError(
        "No PDF converter available. Install LibreOffice (soffice) on Linux, "
        "or Microsoft PowerPoint on Windows."
    )


def _register_chinese_font() -> str:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    workspace_root = Path(__file__).resolve().parent.parent
    candidates = [
        workspace_root / "haitai_week_report" / "fonts" / "NotoSansSC-Regular.otf",
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"),
        Path("/usr/share/fonts/wqy-microhei/wqy-microhei.ttc"),
    ]
    for font_path in candidates:
        if not font_path.exists():
            continue
        name = "ChineseFont"
        try:
            suffix = font_path.suffix.lower()
            if suffix == ".ttc":
                pdfmetrics.registerFont(TTFont(name, str(font_path), subfontIndex=0))
            else:
                pdfmetrics.registerFont(TTFont(name, str(font_path)))
            return name
        except Exception:
            continue
    return "Helvetica"


def add_watermark_to_pdf(input_pdf: Path, output_pdf: Path, text: str) -> None:
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas

    font_name = _register_chinese_font()
    reader = PdfReader(str(input_pdf))
    writer = PdfWriter()

    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)

        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=(width, height))
        can.setFillAlpha(0.12)
        can.setFont(font_name, 36)
        can.saveState()
        can.translate(width * 0.5, height * 0.5)
        can.rotate(35)
        can.setFillColorRGB(0.55, 0.55, 0.55)
        can.drawCentredString(0, 0, text)
        can.restoreState()
        can.save()

        packet.seek(0)
        overlay_reader = PdfReader(packet)
        page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with open(output_pdf, "wb") as f:
        writer.write(f)


def convert_ppt_to_watermarked_pdf(pptx_path: Path, pdf_path: Path) -> None:
    tmp_pdf = pdf_path.with_suffix(".tmp.pdf")
    try:
        pptx_to_pdf(pptx_path, tmp_pdf)
        try:
            add_watermark_to_pdf(tmp_pdf, pdf_path, WATERMARK_TEXT)
        except Exception as watermark_err:
            print(f"Watermark failed, keeping plain PDF: {watermark_err}", file=sys.stderr)
            shutil.copy2(tmp_pdf, pdf_path)
    finally:
        if tmp_pdf.exists():
            tmp_pdf.unlink()


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="私募产品历史业绩月报生成器")
    parser.add_argument(
        "--workspace",
        default=str(WORKSPACE),
        help="工作目录（含净值文件与 report_config.json）",
    )
    parser.add_argument("-o", "--output", help="输出目录")
    parser.add_argument("--end-date", help="报告截止日期 (YYYY-MM-DD)")
    parser.add_argument("--config", help="报告配置 JSON 路径（默认 workspace/report_config.json）")
    parser.add_argument("--theme", choices=["red", "blue"], default="red", help="报告主题色")
    parser.add_argument("--pptx-name", help="输出 PPTX 文件名")
    parser.add_argument("--pdf-name", help="输出 PDF 文件名")
    args = parser.parse_args(argv)

    apply_theme(RED_THEME if args.theme == "red" else BLUE_THEME)

    workspace = Path(args.workspace).resolve()
    config_path = Path(args.config).resolve() if args.config else workspace / "report_config.json"
    config: Optional[dict] = None
    if config_path.is_file():
        config = json.loads(config_path.read_text(encoding="utf-8"))
    if args.end_date:
        if config is None:
            config = {"end_date": args.end_date, "section_order": {}, "products": {}}
        else:
            config["end_date"] = args.end_date

    print("Loading fund data...")
    reports = build_reports(workspace, config)
    if not reports:
        print("错误: 未加载任何产品数据", file=sys.stderr)
        return 1
    print(f"Loaded {len(reports)} products.")

    date_suffix = END_DATE.strftime("%Y%m%d")
    out_dir = Path(args.output).resolve() if args.output else OUTPUT_DIR
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_name = args.pptx_name or f"私募产品历史业绩_{date_suffix}.pptx"
    pdf_name = args.pdf_name or f"私募产品历史业绩_{date_suffix}.pdf"
    pptx_path = out_dir / pptx_name
    pdf_path = out_dir / pdf_name

    print("Generating PPT...")
    generate_ppt(reports, pptx_path)
    print(f"PPT saved: {pptx_path}")

    print("Converting to PDF and adding watermark...")
    pdf_ok = False
    try:
        convert_ppt_to_watermarked_pdf(pptx_path, pdf_path)
        print(f"PDF saved: {pdf_path}")
        pdf_ok = True
    except Exception as exc:
        print(f"PDF conversion failed: {exc}", file=sys.stderr)
        print(
            "PPT was generated successfully. PDF requires LibreOffice (soffice) on Linux "
            "or Microsoft PowerPoint on Windows.",
            file=sys.stderr,
        )

    print("Done.")
    return 0 if pdf_ok else 2


if __name__ == "__main__":
    raise SystemExit(main())
